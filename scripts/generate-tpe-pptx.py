"""
Génère la synthèse Trade-off TPE en PowerPoint.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette couleurs ──────────────────────────────────────────
C_BLUE_DARK   = RGBColor(0x1F, 0x39, 0x64)   # titres slides
C_BLUE_MID    = RGBColor(0x2E, 0x75, 0xB6)   # bandeau titre
C_BLUE_LIGHT  = RGBColor(0xD6, 0xE4, 0xF0)   # lignes paires
C_GREEN       = RGBColor(0x37, 0x86, 0x44)
C_ORANGE      = RGBColor(0xC5, 0x5A, 0x11)
C_RED         = RGBColor(0xC0, 0x00, 0x00)
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_GREY_LIGHT  = RGBColor(0xF2, 0xF2, 0xF2)
C_GREY_TEXT   = RGBColor(0x40, 0x40, 0x40)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # layout vide

# ── Helpers ───────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=12, bold=False, color=C_GREY_TEXT,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def slide_header(slide, title, subtitle=""):
    """Bandeau bleu en haut + titre blanc."""
    add_rect(slide, 0, 0, 13.33, 1.15, fill=C_BLUE_MID)
    add_text(slide, title, 0.3, 0.12, 12, 0.6,
             size=24, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.3, 0.72, 12, 0.4,
                 size=13, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT, italic=True)

def footer(slide, txt="Trade-off TPE — Association | Mai 2026"):
    add_rect(slide, 0, 7.2, 13.33, 0.3, fill=C_BLUE_DARK)
    add_text(slide, txt, 0.3, 7.22, 12, 0.26,
             size=9, color=C_WHITE, align=PP_ALIGN.LEFT)

def icon_cell(slide, icon, label, l, t, w=3.5):
    """Bloc icône + label centré."""
    add_rect(slide, l, t, w, 0.9, fill=C_BLUE_LIGHT)
    add_text(slide, icon,  l+0.1, t+0.05, 0.7, 0.8, size=20, align=PP_ALIGN.CENTER)
    add_text(slide, label, l+0.8, t+0.2,  w-1,  0.5, size=11, bold=True, color=C_BLUE_DARK)

def table_header_row(slide, cols, tops, widths, row_top, row_h=0.38):
    """Ligne d'en-tête bleue foncée."""
    for col, left, w in zip(cols, tops, widths):
        add_rect(slide, left, row_top, w-0.02, row_h, fill=C_BLUE_DARK)
        add_text(slide, col, left+0.07, row_top+0.05, w-0.15, row_h-0.06,
                 size=9.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

def table_row(slide, cells, tops, widths, row_top, row_h=0.38, shade=False):
    bg = C_BLUE_LIGHT if shade else C_WHITE
    for cell, left, w in zip(cells, tops, widths):
        add_rect(slide, left, row_top, w-0.02, row_h, fill=bg,
                 line=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.5))
        # couleur texte selon indicateur
        col = C_GREY_TEXT
        if cell.startswith("✅"): col = C_GREEN
        elif cell.startswith("⚠️"): col = C_ORANGE
        elif cell.startswith("❌"): col = C_RED
        add_text(slide, cell, left+0.08, row_top+0.04, w-0.18, row_h-0.06,
                 size=9, color=col)

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — Titre
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_BLUE_DARK)
add_rect(sl, 0, 2.5, 13.33, 2.5, fill=C_BLUE_MID)

add_text(sl, "Trade-off TPE", 1, 2.65, 11, 1.1,
         size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "QR Invoice  ·  S700  ·  WisePad 3",
         1, 3.65, 11, 0.7, size=22, color=C_WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "Association sportive  |  CA 20 000 €/an  |  Mai 2026",
         1, 6.8, 11, 0.5, size=12, italic=True,
         color=RGBColor(0xBF,0xD7,0xED), align=PP_ALIGN.CENTER)

add_text(sl, "Hypothèses : 15 % carte physique (3 000 €/an) · 5 % online (1 000 €/an)",
         1, 5.8, 11, 0.5, size=11, italic=True,
         color=RGBColor(0xBF,0xD7,0xED), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — Architecture (qui fait quoi)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Architecture", "Qui fait quoi — responsabilités machine")
footer(sl)

cols   = ["", "QR Invoice", "S700", "WisePad 3"]
lefts  = [0.15, 3.1, 6.15, 9.2]
widths = [2.95, 3.05, 3.05, 3.95]
top    = 1.3

table_header_row(sl, cols, lefts, widths, top)

rows = [
    ("App caisse",         "Commande + enregistre",          "Commande + enregistre",        "Commande + relaie les données"),
    ("Matériel / client",  "Téléphone du client (QR)",       "S700 exécute les ordres Stripe","WisePad lit la carte uniquement"),
    ("Stripe serveur",     "Tout : page, paiement, sécurité","Pilote le TPE + encaisse",      "Encaisse (reçoit depuis l'app)"),
    ("Maillon faible",     "Connexion internet client",      "Connexion internet caisse",     "Bluetooth + navigateur Chrome"),
]
for i,(label,c1,c2,c3) in enumerate(rows):
    rt = top + 0.38*(i+1)
    table_row(sl, [label,c1,c2,c3], lefts, widths, rt, shade=(i%2==1))

# Boîtes synthèse
add_rect(sl, 0.15, 3.0, 2.9, 1.2, fill=C_GREY_LIGHT, line=C_BLUE_MID, line_w=Pt(1))
add_rect(sl, 3.1,  3.0, 3.0, 1.2, fill=C_GREY_LIGHT, line=C_BLUE_MID, line_w=Pt(1))
add_rect(sl, 6.15, 3.0, 3.0, 1.2, fill=C_GREY_LIGHT, line=C_BLUE_MID, line_w=Pt(1))
add_rect(sl, 9.2,  3.0, 3.9, 1.2, fill=C_GREY_LIGHT, line=C_ORANGE,   line_w=Pt(1.5))

add_text(sl, "Rôle pivot :", 0.25, 3.05, 2.7, 0.3, size=9, bold=True, color=C_BLUE_DARK)
add_text(sl, "Stripe porte tout.\nL'app ne touche jamais la carte.",
         0.25, 3.35, 2.7, 0.8, size=9, color=C_GREY_TEXT)
add_text(sl, "Rôle pivot :", 3.2, 3.05, 2.8, 0.3, size=9, bold=True, color=C_BLUE_DARK)
add_text(sl, "Stripe pilote le S700\nà distance via Internet.",
         3.2, 3.35, 2.8, 0.8, size=9, color=C_GREY_TEXT)
add_text(sl, "Rôle pivot :", 6.25, 3.05, 2.8, 0.3, size=9, bold=True, color=C_BLUE_DARK)
add_text(sl, "L'app est le câble.\nBluetooth + Chrome requis.",
         6.25, 3.35, 2.8, 0.8, size=9, color=C_GREY_TEXT)
add_text(sl, "⚠️  Risque spécifique :", 9.3, 3.05, 3.6, 0.3, size=9, bold=True, color=C_ORANGE)
add_text(sl, "Panne possible sans panne\nmatérielle (navigateur, OS).",
         9.3, 3.35, 3.6, 0.8, size=9, color=C_ORANGE)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — Coûts & TCO
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Coûts & TCO sur 3 ans", "Hypothèses : 3 000 €/an carte physique · 1 000 €/an online")
footer(sl)

cols   = ["Critère", "QR Invoice", "S700", "WisePad 3"]
lefts  = [0.2, 3.4, 6.6, 9.8]
widths = [3.2, 3.2, 3.2, 3.3]
top    = 1.25

table_header_row(sl, cols, lefts, widths, top)

rows = [
    ("Matériel",                     "0 €",       "349 €",     "79 €"),
    ("Amortissement / an (3 ans)",   "0 €",       "~116 €/an", "~26 €/an"),
    ("Fee Stripe / tx carte",        "1.5% + 0.25€", "1.5% + 0.10€", "1.5% + 0.10€"),
    ("Fees carte annuels (3 000 €)", "75 €/an",   "57 €/an",   "57 €/an"),
    ("Fees online annuels (1 000 €)","25 €/an",   "25 €/an",   "25 €/an"),
    ("TCO 3 ans",                    "300 €",      "595 €",     "325 €"),
    ("Break-even matériel",          "—",          "~19 ans ❌","~4.4 ans ⚠️"),
]
for i,(label,c1,c2,c3) in enumerate(rows):
    rt = top + 0.38*(i+1)
    # mettre en évidence la ligne TCO
    shade = (i%2==1)
    table_row(sl, [label,c1,c2,c3], lefts, widths, rt, shade=shade)

# Note sous le tableau
add_rect(sl, 0.2, 4.35, 12.9, 0.55, fill=RGBColor(0xFF,0xF2,0xCC),
         line=C_ORANGE, line_w=Pt(0.75))
add_text(sl, "💡  La différence de fees entre QR (0.25 €/tx) et card present (0.10 €/tx) ne génère que ~18 €/an d'économies "
             "— insuffisant pour amortir le S700 dans un délai raisonnable.",
         0.35, 4.38, 12.6, 0.5, size=9.5, color=C_ORANGE)

# Graphique barre simplifié (rectangles proportionnels)
add_text(sl, "TCO 3 ans comparé", 0.2, 5.05, 4, 0.35, size=11, bold=True, color=C_BLUE_DARK)

bar_data = [("QR Invoice", 300, C_GREEN), ("WisePad 3", 325, C_BLUE_MID), ("S700", 595, C_ORANGE)]
bar_max = 595
bar_l, bar_t, bar_max_w = 0.2, 5.5, 8.0
for i,(label, val, col) in enumerate(bar_data):
    bt = bar_t + i*0.6
    bw = bar_max_w * val / bar_max
    add_rect(sl, bar_l, bt, bw, 0.42, fill=col)
    add_text(sl, f"{label}  {val} €", bar_l + bw + 0.1, bt+0.07, 3, 0.3,
             size=10, bold=True, color=col)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture logicielle & Testabilité
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Architecture logicielle & Testabilité", "Complexité d'implémentation et de validation")
footer(sl)

cols   = ["Critère", "QR Invoice", "S700", "WisePad 3"]
lefts  = [0.2, 3.4, 6.6, 9.8]
widths = [3.2, 3.2, 3.2, 3.3]
top    = 1.25

table_header_row(sl, cols, lefts, widths, top)

rows = [
    ("Complexité code",    "✅ Faible\nRéutilise infra online\ndéjà en prod",
                           "⚠️ Moyenne\n5 étapes async\n+ nouveau SDK + Lambda",
                           "❌ Élevée\nIdentique S700\n+ relai Bluetooth"),
    ("Testabilité",        "✅ Excellent\n1 HTTP + 1 webhook\nCI / tous navigateurs",
                           "✅ Bon\nSimulateur S700 fidèle\nhappy path couvert",
                           "❌ Partiel\nSimulateur ≠ Bluetooth réel\nvalidation physique obligatoire"),
    ("Développement",      "~2 jours",     "~4 jours\n(code déjà prêt ✅)", "Non retenu"),
    ("État actuel",        "❌ À implémenter", "✅ Code prêt,\ndésactivé",  "❌ Abandonné"),
]
for i,(label,c1,c2,c3) in enumerate(rows):
    rt = top + 0.7*(i+1) - 0.3
    h = 0.68
    table_row(sl, [label,c1,c2,c3], lefts, widths, rt, row_h=h, shade=(i%2==1))

# Note WisePad
add_rect(sl, 0.2, 5.5, 12.9, 0.6, fill=RGBColor(0xFF,0xEB,0xEB),
         line=C_RED, line_w=Pt(0.75))
add_text(sl, "❌  WisePad 3 abandonné : Web Bluetooth API non supporté par Firefox/Safari. "
             "Le simulateur Stripe retourne un reader internet générique — le Bluetooth réel n'est jamais testé.",
         0.35, 5.53, 12.5, 0.55, size=9.5, color=C_RED)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — Sécurité
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Sécurité", "Données carte · Surface d'attaque · Certification PCI")
footer(sl)

cols   = ["Critère", "QR Invoice", "S700", "WisePad 3"]
lefts  = [0.2, 3.4, 6.6, 9.8]
widths = [3.2, 3.2, 3.2, 3.3]
top    = 1.25

table_header_row(sl, cols, lefts, widths, top)

rows = [
    ("Données carte vues par l'app",
     "✅ Jamais\nL'app génère une URL,\nc'est tout",
     "✅ Jamais\nChiffrées dans le S700\nhardware certifié P2PE",
     "✅ Jamais\nChiffrées dans le WisePad\n(mais relai navigateur)"),
    ("Surface d'attaque app",
     "✅ Minimale",
     "✅ Minimale",
     "⚠️ Plus large\nRelai BT → Stripe\ndans le navigateur"),
    ("Certification PCI",
     "✅ 100% déléguée à Stripe",
     "✅ S700 tamper-proof\nPCI P2PE hardware",
     "✅ WisePad certifié\n⚠️ Maillon navigateur\nhors périmètre PCI"),
    ("Verdict",
     "✅ Excellent",
     "✅ Excellent",
     "⚠️ Bon"),
]
for i,(label,c1,c2,c3) in enumerate(rows):
    rt = top + 0.72*(i+1) - 0.3
    table_row(sl, [label,c1,c2,c3], lefts, widths, rt, row_h=0.70, shade=(i%2==1))

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — Robustesse aux pannes
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Robustesse aux pannes", "Population âgée · Bénévoles en caisse")
footer(sl)

cols   = ["Panne", "QR Invoice", "S700", "WisePad 3"]
lefts  = [0.2, 3.9, 6.85, 9.8]
widths = [3.7, 2.95, 2.95, 3.3]
top    = 1.25

table_header_row(sl, cols, lefts, widths, top)

rows = [
    ("Internet caisse coupé",
     "❌  show-stopper",
     "❌  show-stopper",
     "❌  show-stopper"),
    ("Internet / mobile CLIENT absent\n(data épuisé, mauvaise réception, pas de smartphone)",
     "❌  SPÉCIFIQUE\nPaiement impossible",
     "✅  Sans objet",
     "✅  Sans objet"),
    ("Panne navigateur / mise à jour OS",
     "✅",
     "✅",
     "❌  SPÉCIFIQUE\nBluetooth peut casser\nsans avertissement"),
    ("Panne matérielle",
     "✅  Aucun matériel",
     "⚠️  349 € à remplacer",
     "⚠️  79 €, USB-C rechargeable"),
    ("Panne silencieuse inexplicable",
     "✅  Impossible",
     "✅  Toujours détectable",
     "❌  Risque réel (BT + Chrome)"),
]
for i,(label,c1,c2,c3) in enumerate(rows):
    rt = top + 0.5*(i+1)
    h = 0.48
    table_row(sl, [label,c1,c2,c3], lefts, widths, rt, row_h=h, shade=(i%2==1))

# Encadré warning
add_rect(sl, 0.2, 4.0, 12.9, 0.85, fill=RGBColor(0xFF,0xF2,0xCC),
         line=C_ORANGE, line_w=Pt(1.5))
add_text(sl, "⚠️  Point critique pour ce contexte",
         0.4, 4.05, 8, 0.3, size=11, bold=True, color=C_ORANGE)
add_text(sl, "Le QR Invoice exige deux connexions Internet simultanées (caisse + téléphone client). "
             "Pour une population âgée, le taux d'échec est structurel et prévisible : "
             "pas de smartphone, data épuisé, mauvaise réception en salle de sport.",
         0.4, 4.38, 12.5, 0.45, size=9.5, color=C_ORANGE)

# Hiérarchie
add_text(sl, "Hiérarchie robustesse réelle :", 0.2, 5.0, 5, 0.35,
         size=11, bold=True, color=C_BLUE_DARK)
add_text(sl, "S700  ≥  WisePad 3  >>  QR Invoice",
         0.2, 5.38, 8, 0.45, size=16, bold=True, color=C_BLUE_DARK)

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — Tableau de décision final
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Tableau de décision", "Synthèse multi-critères")
footer(sl)

cols   = ["Critère",        "QR Invoice", "S700",       "WisePad 3"]
lefts  = [0.2, 3.7, 6.65, 9.6]
widths = [3.5, 2.95, 2.95, 3.55]
top    = 1.25

table_header_row(sl, cols, lefts, widths, top)

rows = [
    ("Prix / TCO",                  "✅ Gagnant (300 €)",  "❌ Élevé (595 €)",   "⚠️ Correct (325 €)"),
    ("Architecture simple",         "✅ Gagnant",           "⚠️ Moyenne",          "❌ Complexe"),
    ("Testabilité",                  "✅ Gagnant",           "✅ Bon",               "❌ Partiel"),
    ("Sécurité",                    "✅ Excellent",         "✅ Excellent",         "⚠️ Bon"),
    ("Robustesse (pop. âgée)",      "❌ Risque structurel", "✅ Gagnant",           "⚠️ Acceptable"),
    ("UX bénévole en caisse",       "⚠️ Rupture workflow",  "✅ Gagnant (familier)","✅ Familier"),
]
for i,(label,c1,c2,c3) in enumerate(rows):
    rt = top + 0.42*(i+1)
    table_row(sl, [label,c1,c2,c3], lefts, widths, rt, shade=(i%2==1))

# Ligne récap colorée
add_rect(sl, 0.2, 4.0, 12.9, 0.42, fill=C_BLUE_DARK)
add_text(sl, "VERDICT", 0.4, 4.05, 3.2, 0.35, size=11, bold=True, color=C_WHITE)
add_text(sl, "Option économique\nsi clients autonomes",
         3.8, 4.05, 2.7, 0.35, size=9, color=RGBColor(0xBF,0xD7,0xED))
add_text(sl, "✅ Recommandé\nsi population âgée",
         6.75, 4.05, 2.7, 0.35, size=9, bold=True, color=C_WHITE)
add_text(sl, "Abandonné\n(Bluetooth instable)",
         9.7, 4.05, 3.3, 0.35, size=9, color=RGBColor(0xBF,0xD7,0xED))

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — Recommandation & Next Steps
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Recommandation & Prochaines étapes", "")
footer(sl)

# Deux colonnes : population autonome vs âgée
add_rect(sl, 0.2, 1.3, 6.1, 4.5, fill=C_GREY_LIGHT, line=C_BLUE_MID, line_w=Pt(1))
add_rect(sl, 6.5, 1.3, 6.6, 4.5, fill=RGBColor(0xE2,0xEF,0xDA), line=C_GREEN, line_w=Pt(2))

add_text(sl, "Option A — Membres autonomes (smartphone)", 0.35, 1.38, 5.8, 0.4,
         size=11, bold=True, color=C_BLUE_DARK)
add_text(sl,
         "→ QR Invoice\n\n"
         "• TCO 3 ans : 300 € (le moins cher)\n"
         "• 0 € matériel\n"
         "• Architecture simple, infra déjà en place\n"
         "• ~2 jours de développement\n\n"
         "⚠️ Condition : s'assurer que les membres\n"
         "ont tous un smartphone avec data",
         0.35, 1.85, 5.8, 3.8, size=10, color=C_GREY_TEXT)

add_text(sl, "✅ Option B — Population âgée / bénévoles (RECOMMANDÉ)", 6.65, 1.38, 6.2, 0.4,
         size=11, bold=True, color=C_GREEN)
add_text(sl,
         "→ S700\n\n"
         "• Robustesse maximale : zéro dépendance\n"
         "  connexion côté client\n"
         "• UX identique aux TPE bancaires\n"
         "• Code déjà prêt, simulateur fiable\n"
         "• Surcoût réel : ~295 € sur 3 ans\n\n"
         "Prochaines étapes :\n"
         "1. Déployer Lambda stripe-connection-token\n"
         "2. Activer tpe_payment_active: true\n"
         "3. Commander le S700 Stripe (~349 €)",
         6.65, 1.85, 6.2, 3.8, size=10, color=C_GREY_TEXT)

# WisePad 3 — abandonné
add_rect(sl, 0.2, 5.95, 12.9, 0.55, fill=RGBColor(0xF2,0xF2,0xF2),
         line=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.75))
add_text(sl, "❌  WisePad 3 — abandonné définitivement : Bluetooth navigateur instable, "
             "Chrome/Edge uniquement, panne silencieuse impossible à diagnostiquer en caisse.",
         0.35, 5.98, 12.5, 0.5, size=9, italic=True,
         color=RGBColor(0x80,0x80,0x80))

# ── Sauvegarde ────────────────────────────────────────────────
out = r"c:\Users\chrre\Develop\MultipleApps\TPE-Trade-off-Synthese.pptx"
prs.save(out)
print(f"✅ Fichier généré : {out}")
